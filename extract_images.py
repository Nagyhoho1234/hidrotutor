"""
Extract all images from PPTX files in originalMaterials/ and create a catalog.

Outputs:
  - images/ppt/<filename>/  — extracted images organized by source file
  - images/image_catalog.md — markdown catalog mapping images to book chapters
"""

import os
import sys
from pathlib import Path
from collections import defaultdict

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Configuration ──────────────────────────────────────────────────────────
BASE_DIR = Path(r"C:\maidment_hidroGIS")
SRC_DIR = BASE_DIR / "originalMaterials"
OUT_DIR = BASE_DIR / "images" / "ppt"
CATALOG_PATH = BASE_DIR / "images" / "image_catalog.md"

CHAPTER_MAP = {
    "Lecture12018.pptx":              "Ch 1 – Introduction to GIS in Water Resources",
    "MapProj.pptx":                   "Ch 2 – Geodesy and Map Projections",
    "Lecture22018.pptx":              "Ch 4 – GIS Software (ArcGIS Pro)",
    "DataSources2018.pptx":          "Ch 3 – Geospatial Data Sources",
    "SpatialAnalysis.pptx":          "Ch 5-6 – Raster Analysis and Map Algebra",
    "DemWatershedDelineation.pptx":  "Ch 9-11 – Terrain Analysis (DEM & Watersheds)",
    "ExtendedTerrainAnalysis.pptx":  "Ch 9-11 – Terrain Analysis (Extended)",
    "HAND.pptx":                     "Ch 14 – Height Above Nearest Drainage (HAND)",
    "NWMFlood.pptx":                 "Ch 16 – Flood Forecasting (National Water Model)",
    "TFRS.pptx":                     "Ch 16 – Flood Response (Texas Flood Response System)",
    "LIDAR.pptx":                    "Ch 13 – LiDAR Remote Sensing",
    "Programming.pptx":             "Ch 12 – Python Programming for GIS",
    "HydroShare.pptx":              "Ch 3 – Data Sharing (HydroShare)",
    "Demography2018.pptx":          "Ch 17 – Demography and GIS",
}


def emu_to_inches(emu_val):
    """Convert EMU to inches (1 inch = 914400 EMU)."""
    if emu_val is None:
        return None
    return round(emu_val / 914400, 2)


def get_nearby_text(slide, shape):
    """Collect text from shapes near/overlapping the image on the same slide."""
    texts = []
    for s in slide.shapes:
        if s.shape_id == shape.shape_id:
            continue
        if s.has_text_frame:
            txt = s.text_frame.text.strip()
            if txt:
                texts.append(txt)
    return texts


def content_type_to_ext(content_type):
    """Map image content-type to file extension."""
    mapping = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/x-emf": ".emf",
        "image/x-wmf": ".wmf",
        "image/svg+xml": ".svg",
    }
    return mapping.get(content_type, ".bin")


def process_pptx(pptx_path):
    """Process a single PPTX file. Returns list of image records."""
    records = []
    fname = pptx_path.name
    stem = pptx_path.stem
    out_folder = OUT_DIR / stem
    out_folder.mkdir(parents=True, exist_ok=True)

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f"  ERROR opening {fname}: {e}")
        return records

    img_counter = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            is_image = False
            image_blob = None
            content_type = None

            # Direct picture shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                is_image = True
                try:
                    image_blob = shape.image.blob
                    content_type = shape.image.content_type
                except Exception:
                    pass

            # Placeholder with an image
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                try:
                    if hasattr(shape, "image"):
                        image_blob = shape.image.blob
                        content_type = shape.image.content_type
                        is_image = True
                except Exception:
                    pass

            # Group shapes — check children
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            img_counter += 1
                            try:
                                blob = child.image.blob
                                ct = child.image.content_type
                            except Exception:
                                blob = None
                                ct = "unknown"
                            ext = content_type_to_ext(ct) if ct else ".bin"
                            out_name = f"slide{slide_idx:02d}_img{img_counter:03d}{ext}"
                            if blob:
                                out_path = out_folder / out_name
                                out_path.write_bytes(blob)

                            alt_text = ""
                            try:
                                alt_text = child.name or ""
                            except Exception:
                                pass

                            records.append({
                                "file": fname,
                                "slide": slide_idx,
                                "img_num": img_counter,
                                "out_name": out_name,
                                "width_in": emu_to_inches(child.width),
                                "height_in": emu_to_inches(child.height),
                                "left_in": emu_to_inches(child.left),
                                "top_in": emu_to_inches(child.top),
                                "content_type": ct or "unknown",
                                "alt_text": alt_text,
                                "slide_texts": [],
                                "saved": blob is not None,
                            })
                except Exception:
                    pass
                continue

            if not is_image:
                continue

            img_counter += 1
            ext = content_type_to_ext(content_type) if content_type else ".bin"
            out_name = f"slide{slide_idx:02d}_img{img_counter:03d}{ext}"

            if image_blob:
                out_path = out_folder / out_name
                out_path.write_bytes(image_blob)

            # Alt text
            alt_text = ""
            try:
                alt_text = shape.name or ""
            except Exception:
                pass

            # Nearby text on same slide
            slide_texts = get_nearby_text(slide, shape)

            records.append({
                "file": fname,
                "slide": slide_idx,
                "img_num": img_counter,
                "out_name": out_name,
                "width_in": emu_to_inches(shape.width),
                "height_in": emu_to_inches(shape.height),
                "left_in": emu_to_inches(shape.left),
                "top_in": emu_to_inches(shape.top),
                "content_type": content_type or "unknown",
                "alt_text": alt_text,
                "slide_texts": slide_texts,
                "saved": image_blob is not None,
            })

    print(f"  {fname}: {len(prs.slides)} slides, {img_counter} images extracted")
    return records


def build_catalog(all_records):
    """Write image_catalog.md."""
    # Group by source file
    by_file = defaultdict(list)
    for r in all_records:
        by_file[r["file"]].append(r)

    lines = []
    lines.append("# Image Catalog — Extracted from Original PPTX Lectures")
    lines.append("")
    lines.append(f"Generated: 2026-03-25")
    lines.append("")
    lines.append(f"**Total images extracted: {len(all_records)}**")
    lines.append("")

    # Summary table
    lines.append("## Summary by Source File")
    lines.append("")
    lines.append("| Source PPTX | Chapter | Slides with Images | Total Images |")
    lines.append("|---|---|---|---|")
    for fname in sorted(by_file.keys()):
        recs = by_file[fname]
        chapter = CHAPTER_MAP.get(fname, "Unmapped")
        slides_with_imgs = len(set(r["slide"] for r in recs))
        lines.append(f"| {fname} | {chapter} | {slides_with_imgs} | {len(recs)} |")
    lines.append("")

    # Detailed listing per file
    lines.append("---")
    lines.append("")
    lines.append("## Detailed Image Listing")
    lines.append("")

    for fname in sorted(by_file.keys()):
        recs = by_file[fname]
        chapter = CHAPTER_MAP.get(fname, "Unmapped")
        stem = Path(fname).stem
        lines.append(f"### {fname}")
        lines.append(f"**Chapter:** {chapter}")
        lines.append(f"**Output folder:** `images/ppt/{stem}/`")
        lines.append("")
        lines.append("| # | Slide | Filename | Size (W x H in.) | Type | Shape Name / Alt Text | Slide Context |")
        lines.append("|---|---|---|---|---|---|---|")

        for i, r in enumerate(recs, 1):
            w = r["width_in"] or "?"
            h = r["height_in"] or "?"
            ct = r["content_type"].replace("image/", "")
            alt = r["alt_text"].replace("|", "/") if r["alt_text"] else ""
            # Truncate slide texts to keep table readable
            ctx = "; ".join(r["slide_texts"])
            if len(ctx) > 120:
                ctx = ctx[:117] + "..."
            ctx = ctx.replace("|", "/").replace("\n", " ")
            saved = "" if r["saved"] else " (not saved)"
            lines.append(f"| {i} | {r['slide']} | `{r['out_name']}`{saved} | {w} x {h} | {ct} | {alt} | {ctx} |")

        lines.append("")

    CATALOG_PATH.parent.mkdir(parents=True, exist_ok=True)
    CATALOG_PATH.write_text("\n".join(lines), encoding="utf-8")
    print(f"\nCatalog written to {CATALOG_PATH}")


def main():
    print("Scanning PPTX files in originalMaterials/\n")
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    pptx_files = sorted(SRC_DIR.glob("*.pptx"))
    print(f"Found {len(pptx_files)} PPTX files\n")

    all_records = []
    for pf in pptx_files:
        records = process_pptx(pf)
        all_records.extend(records)

    print(f"\nTotal images extracted: {len(all_records)}")
    build_catalog(all_records)


if __name__ == "__main__":
    main()
