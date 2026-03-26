import os, sys, concurrent.futures
from pptx import Presentation

BASE = r"C:\maidment_hidroGIS\originalMaterials"
OUT = os.path.join(BASE, "extracted_text")

FILES = {
    "Lecture12018.pptx": "lecture1.txt",
    "Lecture22018.pptx": "lecture2.txt",
    "MapProj.pptx": "map_projections.txt",
    "DataSources2018.pptx": "data_sources.txt",
    "SpatialAnalysis.pptx": "spatial_analysis.txt",
    "AverageElevation.pptx": "average_elevation.txt",
    "Ex3Overview.pptx": "ex3_overview.txt",
    "Ex3Explanations.pptx": "ex3_explanations.txt",
    "DemWatershedDelineation.pptx": "dem_watershed.txt",
    "ExtendedTerrainAnalysis.pptx": "extended_terrain.txt",
    "FloodConditions.pptx": "flood_conditions.txt",
    "Programming.pptx": "programming.txt",
    "HydroShare.pptx": "hydroshare.txt",
    "Demography2018.pptx": "demography.txt",
    "MidtermReview2018.pptx": "midterm_review.txt",
    "ReviewSlides.pptx": "review_slides.txt",
    "Maidment.pptx": "maidment.txt",
    "FlorenceForecast.pptx": "florence_forecast.txt",
    "TFRS.pptx": "tfrs.txt",
}

MAX_SIZE = 100 * 1024 * 1024  # 100 MB

def extract(pptx_name, txt_name):
    src = os.path.join(BASE, pptx_name)
    dst = os.path.join(OUT, txt_name)
    if not os.path.exists(src):
        return f"MISSING: {pptx_name}"
    sz = os.path.getsize(src)
    if sz > MAX_SIZE:
        return f"SKIPPED (>{MAX_SIZE//1048576}MB): {pptx_name} ({sz//1048576}MB)"
    try:
        prs = Presentation(src)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
            # Also check tables
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct:
                                row_texts.append(ct)
                        if row_texts:
                            lines.append(" | ".join(row_texts))
            lines.append("")  # blank line between slides
        with open(dst, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        slide_count = len(prs.slides)
        return f"OK: {pptx_name} -> {txt_name} ({slide_count} slides)"
    except Exception as e:
        return f"ERROR: {pptx_name}: {e}"

with concurrent.futures.ThreadPoolExecutor(max_workers=8) as pool:
    futures = {pool.submit(extract, k, v): k for k, v in FILES.items()}
    for fut in concurrent.futures.as_completed(futures):
        print(fut.result())
